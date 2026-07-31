from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as C
from pptx.enum.shapes import MSO_SHAPE
import struct
import os

prs = Presentation(); prs.slide_width=In(13.333); prs.slide_height=In(7.5)
W = 13.333

NAVY=C(0x14,0x28,0x4F); GOLD=C(0xB8,0x86,0x0B); SLATE=C(0x46,0x82,0xB4)
RED=C(0xB0,0x3A,0x2E); INK=C(0x21,0x21,0x21); MUTE=C(0x6E,0x6E,0x6E)
BG=C(0xFA,0xFB,0xFC); WHITE=C(0xFF,0xFF,0xFF); RULE=C(0xDC,0xE1,0xE8)
TINT=C(0xEE,0xF3,0xF9); WARM=C(0xFD,0xF6,0xE8)

def tb(slide,x,y,w,h,text,size=12,bold=False,color=INK,align=PP_ALIGN.LEFT,italic=False,space=0):
    t=slide.shapes.add_textbox(In(x),In(y),In(w),In(h)); f=t.text_frame; f.word_wrap=True
    f.margin_left=f.margin_right=f.margin_top=f.margin_bottom=0
    lines = text.split('\n') if isinstance(text,str) else text
    for i,ln in enumerate(lines):
        p=f.paragraphs[0] if i==0 else f.add_paragraph()
        p.text=ln; p.font.size=Pt(size); p.font.bold=bold; p.font.italic=italic
        p.font.color.rgb=color; p.alignment=align; p.font.name='Georgia' if bold and size>=20 else 'Calibri'
        if space: p.space_after=Pt(space)
    return t

def rect(slide,x,y,w,h,fill,line=None,lw=1):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,In(x),In(y),In(w),In(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill; s.shadow.inherit=False
    if line: s.line.color.rgb=line; s.line.width=Pt(lw)
    else: s.line.fill.background()
    return s

def base(title, eyebrow=None):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    rect(s,0,0,W,7.5,BG)
    rect(s,0,0,W,0.98,NAVY)
    rect(s,0,0.98,W,0.045,GOLD)
    if eyebrow: tb(s,0.55,0.14,11,0.25,eyebrow.upper(),8.5,True,GOLD)
        
    tb(s,0.55,0.36,11.8,0.5,title,25,True,WHITE)
    return s

def foot(s,left,right=''):
    rect(s,0.55,6.92,W-1.1,0.012,RULE)
    tb(s,0.55,7.02,8.5,0.3,left,8,False,MUTE)
    if right: tb(s,W-4.6,7.02,4.05,0.3,right,8,False,MUTE,PP_ALIGN.RIGHT)

def band(s,y,text,accent=GOLD,h=0.78,fill=WARM):
    """the 'so what' strip under a chart"""
    rect(s,0.55,y,W-1.1,h,fill)
    rect(s,0.55,y,0.055,h,accent)
    tb(s,0.78,y+0.13,W-1.55,h-0.2,text,10.5,False,INK)

def pic(s,path,x,y,w):
    if os.path.exists(path):
        try: s.shapes.add_picture(path,In(x),In(y),width=In(w)); return True
        except Exception as e: print('  !',path,e)
    else: print('  MISSING',path)
    return False


def png_size(path):
    with open(path,'rb') as f:
        head=f.read(26)
    w,h=struct.unpack('>II', head[16:24])
    return w,h

def fit(s, path, top=1.28, band_h=0.95, bottom=6.82, pad=0.14):
    """Scale image to fill the space between header and the so-what band. Returns band y."""
    band_y = bottom - band_h
    avail_h = band_y - top - pad
    avail_w = W - 1.5
    if not os.path.exists(path):
        print('  MISSING', path); return band_y
    w,h = png_size(path); ar = w/h
    dh = avail_h; dw = dh*ar
    if dw > avail_w:
        dw = avail_w; dh = dw/ar
    x = (W-dw)/2; y = top + (avail_h-dh)/2
    s.shapes.add_picture(path, In(x), In(y), width=In(dw))
    return band_y

def chart_slide(title, eyebrow, img, so_what, foot_l, foot_r, band_h=0.95):
    s = base(title, eyebrow)
    by = fit(s, img, band_h=band_h)
    band(s, by, so_what, h=band_h)
    foot(s, foot_l, foot_r)
    return s

def card(s,x,y,w,h,label,value,note='',vcol=NAVY):
    rect(s,x,y,w,h,WHITE,RULE,1)
    rect(s,x,y,w,0.05,GOLD)
    tb(s,x+0.16,y+0.2,w-0.32,0.26,label.upper(),7.8,True,SLATE)
    tb(s,x+0.16,y+0.5,w-0.32,0.62,value,26,True,vcol)
    if note: tb(s,x+0.16,y+h-0.42,w-0.32,0.3,note,7.8,False,MUTE,italic=True)

def table(s,x,y,w,headers,rows,widths,rh=0.34,fs=9.5):
    cx=x
    for i,hd in enumerate(headers):
        rect(s,cx,y,widths[i],rh,NAVY)
        tb(s,cx+0.1,y+0.075,widths[i]-0.2,rh,hd,8.5,True,WHITE)
        cx+=widths[i]
    for r,row in enumerate(rows):
        cx=x; yy=y+rh+r*rh
        for i,cell in enumerate(row):
            rect(s,cx,yy,widths[i],rh,WHITE if r%2 else TINT)
            col=RED if (isinstance(cell,str) and cell.startswith('!')) else INK
            txt=cell[1:] if (isinstance(cell,str) and cell.startswith('!')) else cell
            tb(s,cx+0.1,yy+0.07,widths[i]-0.2,rh,txt,fs,i==0,col)
            cx+=widths[i]
    return y+rh*(len(rows)+1)

def bullets(s,x,y,w,items,size=11.5,gap=0.34,marker=True):
    for i,it in enumerate(items):
        yy=y+i*gap
        if it=='': continue
        if marker: rect(s,x,yy+0.09,0.075,0.075,GOLD)
        tb(s,x+(0.22 if marker else 0),yy,w-0.22,gap,it,size,False,INK)

D='reports/deck/'; R='reports/'

# ══ 1 TITLE ══
s=prs.slides.add_slide(prs.slide_layouts[6])
rect(s,0,0,W,7.5,NAVY); rect(s,0,0,0.09,7.5,GOLD)
tb(s,0.9,0.85,11,0.3,'BROWN UNIVERSITY INVESTMENT OFFICE',10.5,True,GOLD)
tb(s,0.9,2.25,11.5,1.0,'Private Fund Holdings',46,True,WHITE)
tb(s,0.9,3.35,11.5,0.7,'Data-State Analysis',30,True,GOLD)
rect(s,0.9,4.35,3.2,0.03,C(0x3E,0x55,0x7A))
tb(s,0.9,4.62,11,1.1,['Objective 1 — completeness and cleanliness scorecard across field, fund and manager',
   'DSI Capstone  ·  Phase 1 deliverable'],12.5,False,C(0xC5,0xD0,0xDE),space=6)
for i,(v,l) in enumerate([('52,048','holdings'),('385','funds'),('101','managers'),('8,633','companies'),('329','fields')]):
    x=0.9+i*2.35
    tb(s,x,6.15,2.2,0.45,v,20,True,WHITE)
    tb(s,x,6.62,2.2,0.3,l.upper(),8,True,GOLD)
tb(s,0.9,7.05,11,0.3,'Data period: December 2024 – March 2026  ·  Confidential',8.5,False,C(0x8F,0x9E,0xB4))

# ══ 2 MANDATE ══
s=base('Project Mandate and Scope','Scope')
tb(s,0.55,1.3,12.2,0.4,'The Investment Office receives 400+ fund reports as non-standard PDFs. A third-party vendor flattens them into one CSV. Before that CSV can support any analysis, its trustworthiness has to be measured.',11,False,MUTE)
rect(s,0.55,2.0,6.0,3.5,WHITE,RULE); rect(s,0.55,2.0,6.0,0.05,GOLD)
tb(s,0.8,2.25,5.5,0.3,'MANDATE (PROJECT BRIEF)',9,True,SLATE)
bullets(s,0.8,2.68,5.5,['Identify missing data and its cause',
 'Catch the same company reported under different names',
 'Flag fields that should always be populated but are not',
 'Design repeatable procedures Operations can run'],10.5,0.62)
rect(s,6.78,2.0,6.0,3.5,WHITE,RULE); rect(s,6.78,2.0,6.0,0.05,GOLD)
tb(s,7.03,2.25,5.5,0.3,'THIS DELIVERABLE (OBJECTIVE 1)',9,True,SLATE)
bullets(s,7.03,2.68,5.5,['Completeness scored by field, by fund, by manager',
 'Gap concentration quantified and ranked',
 'Flagged-record extract for GP follow-up',
 'Entity-resolution scope sized for Phase 2'],10.5,0.62)
band(s,5.75,'Success test from the brief: every major field assessed, and findings clear enough for Operations to prioritise remediation without a data scientist in the room.')
foot(s,'Source: Private Fund Holdings — Project Detail, Brown University Investment Office','1')

# ══ 3 METHOD ══
s=base('Methodology','Approach')
steps=[('01','Load and inspect','837 columns profiled; zero-information attributes dropped'),
 ('02','Exclude non-investments','Cash, other assets and liability lines removed from scoring'),
 ('03','Score by field','Fill rates overall and by functional category'),
 ('04','Score by fund','385 vehicles ranked; outliers isolated at ±2σ'),
 ('05','Score by manager','101 GPs ranked; within-GP fund variance measured'),
 ('06','Size entity resolution','Companies appearing under more than one fund counted'),
 ('07','Export','Flagged records and scorecards written for Operations')]
for i,(n,t,d) in enumerate(steps):
    y=1.32+i*0.72
    rect(s,0.55,y,0.62,0.6,NAVY); tb(s,0.55,y+0.15,0.62,0.35,n,13,True,GOLD,PP_ALIGN.CENTER)
    rect(s,1.17,y,7.4,0.6,WHITE,RULE)
    tb(s,1.4,y+0.09,7.0,0.3,t,11.5,True,INK)
    tb(s,1.4,y+0.32,7.0,0.28,d,9.5,False,MUTE)
rect(s,8.85,1.32,3.93,3.55,WHITE,RULE); rect(s,8.85,1.32,3.93,0.05,GOLD)
tb(s,9.08,1.58,3.5,0.3,'SCORING RULE',9,True,SLATE)
tb(s,9.08,2.0,3.5,1.5,['Completeness of a fund =','','non-null mandatory cells','———————————————','real rows × 7 fields'],10,False,INK,space=3)
tb(s,9.08,3.55,3.5,1.2,['Seven mandatory fields:','Deal Status · Current Cost ·','Unrealized Value · Capital Invested','Sector · Geographic Focus · Total Value'],9,False,MUTE,space=2)
rect(s,8.85,5.05,3.93,1.5,WARM); rect(s,8.85,5.05,0.055,1.5,GOLD)
tb(s,9.08,5.25,3.5,1.1,['GUARDRAIL','','Managers with fewer than 20 rows are reported separately. Three holdings and one blank Sector should not read as a reporting failure.'],9,False,INK,space=3)
foot(s,'Scoring is computed on real investment rows only','2')

# ══ 4 PIPELINE ══
s=base('Data Preparation: Column Reduction','Step 1 — Load and Inspect')
pic(s,D+'pipeline.png',1.35,1.34,10.6)
c=[('Columns 100% empty','249','carry no information at all'),
   ('"Unit" metadata twins','258','label columns, not data'),
   ('Columns fully populated','8','of 330 survivors'),
   ('Duplicate rows','0','extract is clean on this axis')]
for i,(l,v,n) in enumerate(c): card(s,0.55+i*3.13,3.55,2.95,1.4,l,v,n)
band(s,5.2,'Two-thirds of the schema is noise. The vendor emits a fixed superset of columns for every fund regardless of strategy, so sparsity here is a schema artefact, not a GP failure — which is why fill rates must be computed after this reduction, not before.')
foot(s,'58,160 rows × 837 columns raw extract','3')

# ══ 5 NON-INVESTMENT ══
s=base('Source Asset Classification','Step 2 — Classify Source Asset')
tb(s,0.55,1.28,12.2,0.36,'The vendor puts every line of a GP schedule into Source Asset. Scoring and entity resolution need different subsets of it.',11,False,MUTE)
table(s,0.55,1.78,12.2,['Class','Example','Rows','Score it','Resolve it'],
 [['holding','Stripe · Valar Co-Invest 13 LP (Qonto)','51,770','Yes','Yes'],
  ['unnamed_aggregate','Other Investments · Seed Investments','278','Yes','!No'],
  ['balance_sheet','Non-Investment Assets — Cash','5,518','!No','!No'],
  ['accounting_entry','Quarterly Unrealized Gain/Loss','577','!No','!No'],
  ['subtotal','Subtotal Private Investments','17','!No','!No']],
 [2.5,4.9,1.5,1.5,1.8],rh=0.38)
rect(s,0.55,4.35,6.0,1.75,WHITE,RULE)
tb(s,0.78,4.56,5.5,1.4,['WHAT THE OLD FILTER MISSED','','It matched only the Non-Investment prefix. 594 rows of journal entries arriving without that prefix stayed in the population — including a singular "Non-Investment Asset - Cash" the regex did not cover.'],9.5,False,INK,space=5)
rect(s,6.78,4.35,6.0,1.75,WARM); rect(s,6.78,4.35,0.055,1.75,GOLD)
tb(s,7.03,4.56,5.5,1.4,['WHY AGGREGATES ARE SPLIT OUT','','"Other Investments" is real exposure and belongs in a score, but it is not an entity. Counting it as a company puts work in the resolution queue that can never be completed.'],9.5,False,INK,space=5)
band(s,6.22,'Matching is anchored to the whole string, not substrings: the book contains real companies called Itz Cash, Keeper Tax Inc. and Gamma Labs (dba Column Tax). SPV and co-invest vehicles are kept as holdings.',h=0.62)
foot(s,'57 regression tests lock this behaviour','4')

# ══ 6 HEADLINES ══
s=base('Portfolio Data Quality Overview','Findings')
m=[('Portfolio completeness','90.6%','mean across 385 funds',NAVY),
   ('Funds at 100%','118','30.6% of the book',NAVY),
   ('Funds below −2σ','19','none above +2σ',RED),
   ('Managers below 85%','14','of 87 scored',RED),
   ('Cross-fund companies','1,107','entity-resolution scope',GOLD),
   ('Flagged records','13,176','ready for GP follow-up',GOLD)]
for i,(l,v,n,col) in enumerate(m):
    card(s,0.55+(i%3)*4.15,1.35+(i//3)*1.75,3.95,1.55,l,v,n,col)
band(s,4.95,'One number hides the story. 90.6% average completeness looks healthy, but the distribution is one-tailed and the variance sits inside managers, not between them — the next eleven slides take that apart.',h=0.95)
foot(s,'All figures computed on 52,048 scorable rows','5')

# ══ 7 FIELD BY CATEGORY ══
chart_slide('Field Completeness by Category','Step 3 — Score by Field',
 D+'field_by_category.png',
 'Identity and classification fields are effectively complete. Financial fields cluster in the low-80s. Returns and ownership fields collapse — and that collapse is the single largest analytical constraint in the dataset.',
 'Gold line = 85% threshold  ·  red bars below 50%','6')

# ══ 8 RETURNS CLIFF ══
s=base('Return Metric Completeness','Finding 1')
tb(s,0.55,1.3,12.2,0.4,'Gross-level return metrics are broadly reported. Net-level metrics are effectively absent.',12,False,MUTE)
table(s,0.55,1.9,7.6,['Return metric','Fill rate','Usable?'],
 [['Gross TVPI','79.1%','Yes'],['RVPI','78.6%','Yes'],['DPI','78.2%','Yes'],
  ['Gross IRR','11.5%','!No'],['Net TVPI','1.3%','!No'],['Net IRR','1.2%','!No']],[3.6,2.2,1.8])
rect(s,8.4,1.9,4.38,2.4,WARM); rect(s,8.4,1.9,0.055,2.4,GOLD)
tb(s,8.65,2.12,3.9,2.0,['WHY THE GAP IS STRUCTURAL','','Net returns are fund-level figures net of fees, carry and expenses. They belong on the capital account statement, not the schedule of investments the vendor parses.','','This is an extraction-scope gap, not a GP reporting failure.'],9.5,False,INK,space=5)
rect(s,0.55,4.5,12.23,1.5,WHITE,RULE)
tb(s,0.8,4.72,11.7,1.2,['CONSEQUENCE FOR THE INVESTMENT OFFICE','','Any net-of-fee performance analysis cannot be sourced from this extract at holding level. Either the capital account reports enter the extraction scope in Phase 2, or net performance stays a fund-accounting output and the holdings CSV is used strictly for gross, position-level work. This should be an explicit decision, not a discovery made mid-analysis.'],10.5,False,INK,space=6)
band(s,6.15,'Recommendation: confirm with the vendor whether capital account statements are in extraction scope before any net-return workstream is committed.',h=0.7)
foot(s,'Dates are similarly thin: Exit Date 6.9%, Initial Investment Date 41.1%','7')

# ══ 9 DEAL STATUS ══
chart_slide('Deal Status Completeness','Finding 2',
 D+'deal_status_mix.png',
 'The brief flags Deal Status as a field that should always carry a value. It is 88.6% filled: 6,639 rows carry no lifecycle marker at all. Those rows cannot be routed to either the current-portfolio view or the realisations view, so they drop silently out of both.',
 'Vocabulary is clean — no typos or inconsistent labels across 52,048 rows','8')

# ══ 10 FIELD DISTRIBUTION ══
s=base('Field-Level Coverage by Fund','Step 3 — Field-Level Distribution')
tb(s,0.55,1.3,12.2,0.4,'A field 90% complete everywhere is a data problem. A field 100% complete in some funds and 0% in others is a policy choice. Counting funds at zero separates the two.',11,False,MUTE)
table(s,0.55,1.95,12.2,['Mandatory field','Funds at 0% fill','Reading'],
 [['Sector','0','Universal — no fund omits it'],
  ['Geographic Focus','0','Universal'],
  ['Unrealized Value','0','Universal'],
  ['Deal Status','3','Near-universal; 3 funds opt out entirely'],
  ['Current Cost','17','Small opt-out cohort'],
  ['Total Value','58','!Structural — 15% of funds never report it'],
  ['Capital Invested','67','!Structural — 17% of funds never report it']],[4.3,2.6,5.3])
band(s,4.95,'Capital Invested and Total Value are not degraded fields — they are absent-by-design for a sixth of the book. Chasing individual blank cells here would waste Operations time; the question is which fund types systematically exclude them and whether that is contractually expected.',h=1.05)
foot(s,'Zero-fill counts computed across 385 funds','9')

# ══ 11 FUND DISTRIBUTION ══
chart_slide('Completeness Distribution: Funds and Managers','Step 4 — Score by Fund',
 R+'completeness_by_dimension.png',
 'Both distributions are left-skewed against a hard ceiling at 100%. The mass sits above the 85% line, which means remediation is a tail problem — a small number of named vehicles, not a portfolio-wide programme.',
 'Left: 385 funds  ·  Right: 87 managers with at least 20 rows','10')

# ══ 12 FUND BANDS ══
s=base('Fund Completeness Ranking','Step 4 — Score by Fund')
table(s,0.55,1.35,7.2,['Completeness band','Funds','Share','Action'],
 [['100%','118','30.6%','None'],['95–99%','92','23.9%','None'],['85–94%','74','19.2%','Monitor'],
  ['75–84%','63','16.4%','Review'],['50–74%','35','9.1%','!Escalate'],['Below 50%','3','0.8%','!Escalate']],
 [2.7,1.2,1.2,2.1])
rect(s,8.0,1.35,4.78,2.4,WHITE,RULE); rect(s,8.0,1.35,4.78,0.05,GOLD)
tb(s,8.23,1.6,4.3,0.3,'THE THREE WORST VEHICLES',9,True,SLATE)
tb(s,8.23,2.0,4.3,1.6,['Ab3ea87 — 36.5%   (z = −4.05)','Ae7bb35 — 42.3%   (z = −3.62)','A80d406 — 47.6%   (z = −3.22)','','Each sits more than three standard deviations below the mean.'],10,False,INK,space=6)
rect(s,8.0,3.95,4.78,1.55,TINT)
tb(s,8.23,4.15,4.3,1.2,['73% of funds already clear the 85% bar. The remediation target is the 38 funds below 75% — a list short enough to work through by name in a single quarter.'],10,False,INK)
band(s,5.75,'Framing for Operations: this is not a data-quality crisis. It is a concentrated tail of 38 named vehicles, of which 3 are severe.')
foot(s,'Bands computed on mandatory-field completeness','11')

# ══ 13 OUTLIERS ══
chart_slide('Outlier Fund Analysis','Step 4 — Outlier Detection',
 D+'outliers.png',
 'Nineteen funds fall below minus two sigma; none rise above plus two. That asymmetry is mechanical — completeness is capped at 100%, so excellence cannot be an outlier while failure is unbounded. Nine funds cluster at exactly 57.1%, which is 4 of 7 mandatory fields: the same four, suggesting a shared template rather than nine independent lapses.',
 'Mean 90.6%  ·  sigma 13.3%','12', band_h=1.05)

# ══ 14 SIZE ══
s=base('Fund Size vs. Data Quality','Step 4 — Size versus Quality')
pic(s,D+'size_quartile.png',0.62,1.34,6.05)
pic(s,R+'fund_size_vs_quality.png',7.0,1.42,5.75)
band(s,4.62,'Mean completeness climbs monotonically from 77.0% to 97.2% across size quartiles, and dispersion collapses from σ 15.5% to 5.0%. But correlation is only 0.25 and small funds span the full range — several report at 100%. Size proxies for back-office maturity; it does not determine outcomes, so it cannot be used to excuse a small manager.',h=1.15)
foot(s,'Quartiles by summed Total Value  ·  n = 385','13')

# ══ 15 MANAGER SCORECARD ══
s=base('Investment Manager Scorecard','Step 5 — Score by Manager')
m=[('Managers scored','87','≥20 rows',NAVY),('Mean','89.8%','',NAVY),
   ('Median','95.1%','well above mean',NAVY),('At 100%','10','',NAVY),
   ('Below 85%','14','16.1% of GPs',RED),('Worst','39.9%','manager A5d1a01',RED)]
for i,(l,v,n,col) in enumerate(m): card(s,0.55+(i%3)*4.15,1.35+(i//3)*1.6,3.95,1.42,l,v,n,col)
table(s,0.55,4.7,12.2,['Least complete managers','Rows','Funds','Completeness','Worst field'],
 [['A5d1a01','169','2','39.9%','Deal Status 18.9%'],
  ['A99a662','54','3','56.3%','Current Cost 0%'],
  ['Aed3782','226','1','57.1%','Current Cost 0%'],
  ['A9fbabf','173','3','62.3%','Capital Invested 0%']],[3.6,1.5,1.5,2.4,3.2],rh=0.32)
band(s,6.25,'Median far above mean confirms a small tail dragging the average. Seven of the bottom ten share the same worst field: Capital Invested.',h=0.62)
foot(s,'14 managers with fewer than 20 rows excluded as low-evidence','14')

# ══ 16 HEATMAP ══
chart_slide('Manager Field Completeness Heatmap','Step 5 — Manager x Field',
 R+'manager_field_heatmap.png',
 'The gaps are vertical, not diffuse. Sector and Geographic Focus stay green across every weak manager; Capital Invested is red almost everywhere. A manager at 60% overall is not sloppy across the board — they are missing one or two specific fields, which makes the follow-up request a single sentence rather than a data-quality lecture.',
 '20 least complete managers  ·  red = low fill','15', band_h=1.05)

# ══ 17 SPREAD ══
chart_slide('Within-Manager Fund Variance','Step 5 — Within-Manager Variance',
 D+'manager_spread.png',
 'Every manager shown has at least one fund at 100% and another below 72%. Same GP, same reporting quarter, opposite results. The capability plainly exists, so these are fund-level breakdowns to raise about a named vehicle — not manager-level policy conversations. Managers with near-zero spread are the opposite case and need the policy conversation instead.',
 'Managers with more than one fund  ·  ranked by spread','16', band_h=1.05)

# ══ 18 INVESTMENT TYPE ══
chart_slide('Completeness by Investment Type','Step 3 — Cross-Cut',
 D+'investment_type.png',
 'Equity-style holdings clear the bar comfortably. Real Estate at 83.4% and Natural Resources at 57.1% sit below it — and both are being measured against a mandatory field set built for equity.',
 'n shown per type  ·  gold line = 85%','17', band_h=0.82)

# ══ 19 REAL ASSETS ══
s=base('Real Assets: Field Definition Mismatch','Finding 3')
table(s,0.55,1.35,12.2,['Mandatory field','Equity meaning','Real-asset reality'],
 [['Deal Status','Current / Exited / Written off','A property is held, leased or listed for sale — no exit event'],
  ['Current Cost','Entry cost of the position','Book value plus capitalised improvements'],
  ['Unrealized Value','Mark versus cost','Third-party appraisal on its own cycle'],
  ['Capital Invested','Staged rounds','One purchase, then ongoing capex'],
  ['Sector','GICS-style classification','Office / retail / industrial taxonomy']],[2.6,3.9,5.7],rh=0.42)
rect(s,0.55,4.05,6.0,1.65,WHITE,RULE)
tb(s,0.78,4.28,5.5,1.3,['NATURAL RESOURCES — 6 ROWS','','Lifecycle is exploration → development → production → decline. None of it maps to exit or write-down. Sample is too small to score meaningfully in any case.'],9.5,False,INK,space=5)
rect(s,6.78,4.05,6.0,1.65,WARM); rect(s,6.78,4.05,0.055,1.65,GOLD)
tb(s,7.03,4.28,5.5,1.3,['WHAT THIS IS NOT','','This is not evidence that real-asset GPs report poorly. It is evidence that the yardstick was built for a different asset class.'],9.5,False,INK,space=5)
band(s,5.95,'Recommendation: define an asset-class-specific mandatory set before these managers appear on any remediation list. Scoring them against the current set produces a false positive.',h=0.72)
foot(s,'Real Estate 991 rows  ·  Natural Resources 6 rows','18')

# ══ 20 CONCENTRATION ══
chart_slide('Portfolio Concentration: Sectors and Companies','Step 6 — Cross-Cut',
 R+'portfolio_concentration.png',
 'Software and Application Software alone account for 32.5% of holdings; the top five sectors reach 44.6%. Concentration this high means a single sector-taxonomy error propagates across a third of the book.',
 'Top 15 sectors  ·  top 10 companies by row count','19', band_h=0.85)

# ══ 21 HYGIENE ══
chart_slide('Source Asset Data Hygiene','Finding 4 — Resolved',
 D+'hygiene_artifacts.png',
 '"Quarterly Unrealized Gain/Loss" appeared 334 times across 64 funds and "ASC 740-10 Accrual" 145 times across 26, both sitting in Source Asset as though they were portfolio companies. The classifier now removes 594 such rows, and both entries drop out of the entity resolution queue where they had been the two largest items.',
 'Detected in the top-20 company frequency table  ·  now excluded','20', band_h=1.15)

# ══ 22 GEOGRAPHY ══
chart_slide('Geographic Distribution of Holdings','Step 6 — Cross-Cut',
 D+'geography.png',
 'Reported completeness reaches 100% in Africa, Peru, Estonia and Nigeria — but each rests on a handful of rows. Small-n geographies should be read as noise, not as reporting excellence.',
 'Top 10 cover 90.4% of holdings','21', band_h=0.85)

# ══ 23 ENTITY RESOLUTION ══
s=base('Entity Resolution Scope','Step 6 — Entity Resolution')
for i,(l,v,n) in enumerate([('Resolvable names','8,633','after classification'),
   ('Held by more than one fund','1,107','12.8% of names'),
   ('Mean appearances per name','6.0','across quarters and funds')]):
    card(s,0.55+i*4.15,1.35,3.95,1.55,l,v,n)
rect(s,0.55,3.2,12.23,1.9,WHITE,RULE); rect(s,0.55,3.2,12.23,0.05,GOLD)
tb(s,0.8,3.45,11.7,1.5,['WHY THE COUNT OVERSTATES THE COMPANY UNIVERSE','',
 'Two funds holding the same business can write it two ways. "Oomnitz" on one schedule of realised gains and "Oomnitza, Inc." on another schedule of investments are the same company — the extract treats them as two.','',
 'The 1,107 multi-fund names are the working set: these are the only ones where a naming mismatch can split a single company into several, and the only ones where cross-fund valuation comparison is possible at all. Twenty-seven names were removed from this queue by the classifier, led by the two journal entries on the previous slide.'],10.5,False,INK,space=5)
band(s,5.35,'Phase 2 dependency: deal-status inference cross-references the schedule of investments against the schedule of realised gains. That join is by company name, so canonical mapping has to land before inference can run. Fuzzy matching may propose candidates; it may not decide them.',h=1.05)
foot(s,'Explicit aliases, not automatic fuzzy acceptance','22')

# ══ 24 FINDINGS MATRIX ══
s=base('Summary of Findings','Synthesis')
table(s,0.55,1.35,12.2,['#','Finding','Scale','Severity','Owner'],
 [['1','Capital Invested absent by design','67 funds (17.4%)','!High','Vendor + GP'],
  ['2','Net-return metrics unusable','Net IRR 1.2%','!High','Vendor scope'],
  ['3','Managers below 85%','14 of 87','!High','Operations'],
  ['4','Entity resolution unresolved','1,107 names','Medium','Phase 2'],
  ['5','Accounting entries as companies','594 rows','Resolved','Fixed'],
  ['6','Deal Status blank','6,639 rows','Medium','GP follow-up'],
  ['7','Real assets mis-scored','997 rows','Low','Redefine metric'],
  ['8','Severe fund outliers','3 funds below 48%','!High','Operations']],
 [0.6,4.5,2.5,2.0,2.6],rh=0.38)
band(s,5.05,'Two of the top three findings are not GP failures — they are extraction-scope and field-definition questions that belong with the vendor. Sequencing those first avoids sending Operations to chase data the GPs already sent.',h=1.0)
foot(s,'Severity weighted by rows affected and analytical impact','23')

# ══ 25 RECOMMENDATIONS ══
s=base('Recommendations','Actions')
recs=[('01','Resolve Capital Invested and Total Value with the vendor',
   'Confirm whether the 67 and 58 zero-fill funds are contractual exclusions or extraction misses. Nothing else should be chased until this is settled.','Vendor','Immediate'),
 ('02','Decide the net-return question explicitly',
   'Bring capital account statements into scope, or record that net performance is out of scope for this extract.','Vendor','Immediate'),
 ('03','Tier manager outreach',
   'Tier 1: 3 GPs below 60%. Tier 2: 11 GPs at 60–85%. Tier 3: high-spread GPs, raised per fund rather than per relationship.','Operations','Q1'),
 ('04','Define an asset-class mandatory set',
   'Real Estate and Natural Resources scored against their own field list before either appears on a remediation list.','Analytics','Q1'),
 ('05','Extend the non-investment filter  ✓ DONE',
   'Three-way classifier shipped with 57 regression tests. 594 rows reclassified; 27 names left the resolution queue.','Analytics','Complete'),
 ('06','Stand up canonical naming',
   'Begin with the 1,107 multi-fund names. Explicit aliases, reviewed, not auto-accepted.','Analytics','Q1–Q2'),
 ('07','Re-run quarterly',
   'The notebook is parameterised. Track band migration rather than the headline average.','Operations','Ongoing')]
for i,(n,t,d,o,w) in enumerate(recs):
    y=1.32+i*0.76
    rect(s,0.55,y,0.5,0.66,NAVY); tb(s,0.55,y+0.19,0.5,0.3,n,11,True,GOLD,PP_ALIGN.CENTER)
    rect(s,1.05,y,9.15,0.66,WHITE,RULE)
    tb(s,1.28,y+0.08,8.7,0.28,t,10.5,True,INK)
    tb(s,1.28,y+0.32,8.7,0.3,d,8.5,False,MUTE)
    rect(s,10.2,y,1.35,0.66,TINT); tb(s,10.2,y+0.22,1.35,0.3,o,8.5,True,SLATE,PP_ALIGN.CENTER)
    rect(s,11.55,y,1.23,0.66,WARM); tb(s,11.55,y+0.22,1.23,0.3,w,8.5,True,INK,PP_ALIGN.CENTER)
foot(s,'Sequenced so vendor questions resolve before GP outreach begins','24')

# ══ 26 PHASE 2 ══
s=base('Phase 1 to Phase 2 Roadmap','Roadmap')
rect(s,0.55,1.35,6.0,4.1,WHITE,RULE); rect(s,0.55,1.35,6.0,0.05,SLATE)
tb(s,0.8,1.62,5.5,0.3,'PHASE 1 — DESCRIPTIVE   ✓ COMPLETE',10,True,SLATE)
tb(s,0.8,1.98,5.5,0.35,'"What state is the data in?"',12,True,INK)
bullets(s,0.8,2.5,5.5,['Completeness scored on three axes','Gap concentration quantified','13,176 records flagged','Entity-resolution scope sized','Eight exports written for Operations'],10,0.44)
rect(s,0.8,4.85,5.5,0.4,TINT); tb(s,0.95,4.94,5.2,0.3,'Cannot answer: is the data correct?',9.5,True,MUTE)
rect(s,6.78,1.35,6.0,4.1,WHITE,RULE); rect(s,6.78,1.35,6.0,0.05,GOLD)
tb(s,7.03,1.62,5.5,0.3,'PHASE 2 — VALIDATIVE   NEXT',10,True,GOLD)
tb(s,7.03,1.98,5.5,0.35,'"Is the data correct?"',12,True,INK)
bullets(s,7.03,2.5,5.5,['Parse schedule of investments from source PDFs','Reconcile cost and fair value to fund totals','Infer Deal Status across schedules','Compare vendor CSV against statements','Attach page and bounding box to every value'],10,0.44)
rect(s,7.03,4.85,5.5,0.4,WARM); tb(s,7.18,4.94,5.2,0.3,'Blocked on: canonical entity naming',9.5,True,INK)
band(s,5.7,'The two phases answer different questions. Completeness says a value is present; validation says it is right. A fund at 100% completeness carrying wrong numbers still scores 100% here — which is precisely why Phase 2 exists.',h=1.0)
foot(s,'Sample statement parsed and reconciled to the dollar','25')

# ══ 27 APPENDIX ══
s=base('Supporting Data and Validation','Appendix')
table(s,0.55,1.35,7.6,['File','Contents'],
 [['data_state_by_field.csv','Fill rate, all 329 value columns'],
  ['data_state_by_fund.csv','Completeness, all 385 funds'],
  ['data_state_by_manager.csv','Completeness, rows, funds — 101 GPs'],
  ['data_state_by_manager_field.csv','Field-level fill per manager'],
  ['flagged_missing_fields.csv','13,176 records for GP follow-up'],
  ['source_asset_classification.csv','Row and name counts per class'],
  ['excluded_non_investment_rows.csv','6,112 rows removed, with reason'],
  ['source_asset_review_queue.csv','4 names held for human decision']],[3.3,4.3],rh=0.36)
rect(s,8.4,1.35,4.38,3.0,WHITE,RULE); rect(s,8.4,1.35,4.38,0.05,GOLD)
tb(s,8.63,1.6,3.9,0.3,'ASSERTIONS PASSED',9,True,SLATE)
bullets(s,8.63,1.98,3.9,['Fund completeness within [0,1]','Manager completeness within [0,1]','No duplicate rows','Unit columns stripped','Flagged ⊆ all rows'],9.5,0.42)
rect(s,0.55,4.05,12.23,1.5,TINT)
tb(s,0.8,4.28,11.7,1.1,['REPRODUCIBILITY','','The notebook runs end to end from the raw extract with no manual steps. Thresholds — the 85% bar, the 20-row manager floor, the mandatory field list — are parameters at the top of the file, so a change in Investment Office policy is a one-line edit and a re-run, not a rewrite.'],10,False,INK,space=5)
band(s,5.75,'Open question for the Investment Office: is 85% the right bar, and should Sector carry equal weight to Current Cost in the mandatory set? Both are currently assumptions.',h=0.72)
foot(s,'Notebook: (phase1)holdings_data_state_analysis.ipynb  ·  62 cells','26')

# ══ 28 CLOSE ══
s=prs.slides.add_slide(prs.slide_layouts[6])
rect(s,0,0,W,7.5,NAVY); rect(s,0,0,0.09,7.5,GOLD)
tb(s,0.9,2.5,11,0.9,'Questions',40,True,WHITE)
rect(s,0.9,3.6,3.2,0.03,C(0x3E,0x55,0x7A))
tb(s,0.9,3.9,11,1.4,['Phase 1 data-state analysis complete.','Phase 2 validation blocked only on canonical entity naming.'],13,False,C(0xC5,0xD0,0xDE),space=8)
tb(s,0.9,6.9,11,0.3,'Brown University Investment Office  ·  Confidential',8.5,False,C(0x8F,0x9E,0xB4))

out='Private_Fund_Holdings_Data_State_Deck.pptx'
prs.save(out)
print(f'\n{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}')
